"""Vision MCP server — 给无视觉主模型（如 DeepSeek）补上读图能力。

把图片（本地路径 / URL / base64）发给火山 Agent Plan 的 doubao-seed-evolving
视觉模型，返回文字描述。走 OpenAI 兼容端点 /api/plan/v3/chat/completions。

配置（环境变量）：
    ARK_AUTH_TOKEN  必填，火山 Agent Plan AUTH_TOKEN（ark- 开头）
    ARK_MODEL       选填，默认 doubao-seed-evolving
    ARK_BASE_URL    选填，默认 https://ark.cn-beijing.volces.com/api/plan/v3

运行：python vision_mcp.py
"""
import base64
import glob
import json
import os
import re
import shutil
import socket
import subprocess
import sys
import tempfile
import time
from typing import Any
from urllib.request import urlopen, Request

from mcp.server.fastmcp import FastMCP

ARK_BASE = os.environ.get("ARK_BASE_URL", "https://ark.cn-beijing.volces.com/api/plan/v3")
MODEL = os.environ.get("ARK_MODEL", "doubao-seed-evolving")
MAX_IMAGE_BYTES = 15 * 1024 * 1024

# Edge 常见安装路径，按顺序探测
_EDGE_CANDIDATES = [
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
]


def _find_edge() -> str:
    for p in _EDGE_CANDIDATES:
        if os.path.isfile(p):
            return p
    found = shutil.which("msedge")
    if found:
        return found
    raise RuntimeError("未找到 msedge.exe，请设置环境变量 EDGE_PATH 指向 Edge 可执行文件")

mcp = FastMCP("vision-mcp")


def _free_port() -> int:
    s = socket.socket()
    s.bind(("127.0.0.1", 0))
    p = s.getsockname()[1]
    s.close()
    return p


def _recv_msg(ws):
    _op, pl = ws.recv_data()
    return json.loads(pl) if pl else {}


def _cdp_full_page(url: str, width: int, height: int) -> str:
    """通过 CDP 拿完整页面高度后整页截图，返回 PNG 临时文件路径。"""
    edge = _find_edge()
    port = _free_port()
    profile = tempfile.mkdtemp(prefix="edge_cdp_")
    proc = subprocess.Popen(
        [edge, f"--remote-debugging-port={port}", f"--user-data-dir={profile}",
         "--headless=new", "--disable-gpu", "--hide-scrollbars",
         "--remote-allow-origins=*", url],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
    )
    fd, png = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    try:
        target = None
        for _ in range(40):
            try:
                t = urlopen(f"http://127.0.0.1:{port}/json", timeout=1).read()
                targets = json.loads(t)
                target = next((x for x in targets if x.get("type") == "page"), None)
                if target:
                    break
            except Exception:
                pass
            time.sleep(0.25)
        if not target:
            raise RuntimeError("Edge CDP 无 page target")

        import websocket
        ws = websocket.create_connection(target["webSocketDebuggerUrl"], timeout=30)

        def cmd(method, params=None, _id=[0]):
            _id[0] += 1
            ws.send(json.dumps({"id": _id[0], "method": method, "params": params or {}}))
            while True:
                m = _recv_msg(ws)
                if m.get("id") == _id[0]:
                    if "error" in m:
                        raise RuntimeError(m["error"])
                    return m.get("result", {})

        cmd("Page.enable")
        time.sleep(0.5)
        metrics = cmd("Page.getLayoutMetrics")
        css_h = int(metrics["cssContentSize"]["height"])
        cmd("Emulation.setDeviceMetricsOverride", {
            "width": width, "height": max(css_h, height),
            "deviceScaleFactor": 1, "mobile": False,
        })
        time.sleep(0.3)
        shot = cmd("Page.captureScreenshot", {"format": "png"})
        with open(png, "wb") as f:
            f.write(base64.b64decode(shot["data"]))
        cmd("Browser.close")
        if not os.path.isfile(png) or os.path.getsize(png) == 0:
            raise RuntimeError("CDP 截图失败（文件为空）")
        return png
    except Exception:
        if os.path.exists(png):
            os.remove(png)
        raise
    finally:
        proc.terminate()
        try:
            proc.wait(timeout=5)
        except Exception:
            proc.kill()
        try:
            shutil.rmtree(profile, ignore_errors=True)
        except Exception:
            pass


def _html_to_png(target: str, width: int = 1280, height: int = 900) -> str:
    """把 HTML（URL / 本地文件 / data URL）渲染成 PNG，返回临时文件路径。

    优先用 CDP 拿完整页面高度整页截图；CDP 不可用时回退到 --screenshot 单屏。
    """
    try:
        return _cdp_full_page(target, width, height)
    except Exception as e:
        try:
            return _html_to_png_fallback(target, width, height)
        except Exception as e2:
            raise RuntimeError(f"CDP 截图失败({e})，回退也失败({e2})")
    finally:
        # CDP 的 finally 已 terminate 并等待；这里再补偿等待，确保 Edge 文件句柄释放
        time.sleep(1.5)


def _html_to_png_fallback(target: str, width: int, height: int) -> str:
    """回退：Edge --screenshot 单屏截图。"""
    edge = _find_edge()
    fd, png = tempfile.mkstemp(suffix=".png")
    os.close(fd)
    try:
        args = [
            edge,
            "--headless=new",
            "--disable-gpu",
            "--hide-scrollbars",
            f"--screenshot={png}",
            f"--window-size={width},{height}",
        ]
        if target.startswith(("http://", "https://", "data:")):
            args.append(target)
        elif os.path.isfile(target):
            args.append("file:///" + os.path.abspath(target).replace("\\", "/"))
        else:
            raise ValueError(f"无法识别的 HTML 来源: {target[:60]}…")
        subprocess.run(args, check=False, timeout=60,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        if not os.path.isfile(png) or os.path.getsize(png) == 0:
            raise RuntimeError("Edge 截图失败（未生成文件）")
        return png
    except Exception:
        if os.path.exists(png):
            os.remove(png)
        raise


def _read_image_bytes(source: str) -> tuple[bytes, str, str]:
    """从本地路径 / URL / data URL / base64 读取图片，返回 (bytes, mime, 来源描述)。"""
    if source.startswith("data:"):
        m = re.match(r"data:([^;]+);base64,(.+)", source, re.S)
        if not m:
            raise ValueError("无法解析 data URL")
        return base64.b64decode(m.group(2)), m.group(1), "data-url"
    if re.match(r"^https?://", source):
        req = Request(source, headers={"User-Agent": "Mozilla/5.0"})
        with urlopen(req, timeout=30) as r:
            data = r.read()
        return data, r.headers.get_content_type() or "image/jpeg", source
    if len(source) > 512 and re.match(r"^[A-Za-z0-9+/=\s]+$", source):
        return base64.b64decode(source), "image/png", "base64"
    if os.path.isfile(source):
        with open(source, "rb") as f:
            data = f.read()
        return data, "image/jpeg", source
    raise ValueError(f"无法识别的图片来源: {source[:60]}…")


def _thinking_payload(payload: dict, thinking: bool | None) -> dict:
    """按 thinking 设置注入火山 thinking 字段。None = 不注入（遵循模型默认）。"""
    if thinking is None:
        return payload
    payload["thinking"] = {"type": "enabled" if thinking else "disabled"}
    return payload


def _chat(
    messages: list[dict],
    max_tokens: int = 2048,
    timeout: float = 90,
    thinking: bool | None = None,
) -> dict:
    """调用火山 Agent Plan /chat/completions，返回 assistant content 文本。

    429 / 5xx 自动重试（最多 3 次，指数退避），应对火山端临时过载。
    max_tokens / timeout 可按需调大（长描述、逐字转录场景火山端响应更慢）。
    thinking: None=遵循模型默认；True/False=强制开/关深度思考。
    """
    token = os.environ.get("ARK_AUTH_TOKEN")
    if not token:
        raise RuntimeError("缺少环境变量 ARK_AUTH_TOKEN（火山 Agent Plan AUTH_TOKEN）")
    payload = {"model": MODEL, "messages": messages, "max_tokens": max_tokens}
    payload = _thinking_payload(payload, thinking)
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
    }

    def _post() -> tuple[int, bytes]:
        try:
            from curl_cffi import requests
        except ImportError:
            from urllib.request import Request as Req, urlopen as open_

            body = json.dumps(payload).encode()
            req = Req(f"{ARK_BASE}/chat/completions", data=body, headers=headers)
            with open_(req, timeout=timeout) as r:
                return getattr(r, "status", 200), r.read()
        else:
            resp = requests.post(
                f"{ARK_BASE}/chat/completions", json=payload, headers=headers, timeout=timeout
            )
            return resp.status_code, resp.content

    for attempt in range(3):
        status, data = _post()
        if status == 200:
            break
        if status in (429, 500, 502, 503, 504) and attempt < 2:
            time.sleep(3 * (attempt + 1))
            continue
        raise RuntimeError(f"火山 API 返回 HTTP {status}: {data.decode('utf-8', 'replace')[:500]}")

    parsed = json.loads(data)
    try:
        return parsed["choices"][0]["message"]["content"]
    except (KeyError, IndexError):
        raise RuntimeError(f"无法解析响应: {data.decode('utf-8', 'replace')[:500]}")


def _analyze(
    image: str,
    prompt: str,
    max_tokens: int = 2048,
    timeout: float = 90,
    thinking: bool | None = None,
) -> str:
    data, mime, _ = _read_image_bytes(image)
    if len(data) > MAX_IMAGE_BYTES:
        raise ValueError(f"图片过大（{len(data)} bytes），上限 {MAX_IMAGE_BYTES}")
    b64 = base64.b64encode(data).decode()
    content: list[dict[str, Any]] = [
        {"type": "text", "text": prompt},
        {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
    ]
    return _chat(
        [{"role": "user", "content": content}],
        max_tokens=max_tokens, timeout=timeout, thinking=thinking,
    )


_DETAILED_TEMPLATE = """我是DeepSeek，没有多模态，请你作为我的眼睛返回图片的信息：
1. 内容：逐字转录所有文字、数字、表格（不要概括、不要改写），保留精确数值。
2. 布局：分区位置、层级、视觉重点，可给出坐标。
3. 图表：坐标轴、刻度、图例、色标数值。
4. 如果有显示异常（损坏/遮挡/乱码），请指出。"""


def _analyze_detailed(
    image: str,
    prompt: str = "",
    thinking: bool | None = None,
) -> str:
    """对图片做完整描述：逐字转录 + 布局/坐标/图表细节 + 异常检测。

    单轮调用（实测查漏补缺轮多数情况返回"无"，已去掉以省时）；
    用更大的 max_tokens/timeout（逐字转录时火山端响应更慢）。
    """
    first_prompt = (
        f"{_DETAILED_TEMPLATE}\n\n补充关注点：{prompt}" if prompt
        else _DETAILED_TEMPLATE
    )
    return _analyze(image, first_prompt, max_tokens=4096, timeout=240, thinking=thinking)


def _resolve_thinking(detail: bool, thinking: bool | None) -> bool | None:
    """按场景决定思考默认值，thinking 显式传入时覆盖。

    默认一律关思考（实测 detail+思考开会卡死 210s+）；需要深度推理时
    显式传 thinking=True。
    """
    if thinking is not None:
        return thinking
    return False


# ---------------------------------------------------------------------------
# 会话日志定位：从 Claude Code 当前会话的 JSONL 日志里取出用户粘贴的图片。
# MCP server 是 Claude Code 的子进程（stdio），会继承 CLAUDE_CODE_SESSION_ID
# 和 CLAUDE_CODE_EXECPATH，可据此定位到 ~/.claude/projects/<项目>/<会话>.jsonl。
# 用户粘贴的图片以 {"type":"image","source":{"type":"base64","media_type":...,"data":...}}
# 形式出现在 type=="user" 的消息里，取最后一条即"最近粘贴的那张图"。
# ---------------------------------------------------------------------------


def _claude_projects_dir() -> str:
    """返回 ~/.claude/projects 目录。"""
    return os.path.join(os.path.expanduser("~"), ".claude", "projects")


def _extract_pasted_image(line: str) -> tuple[bytes, str] | None:
    """从一行 JSONL 中提取用户消息里粘贴的图片，返回 (bytes, mime)。"""
    try:
        obj = json.loads(line)
    except json.JSONDecodeError:
        return None
    if obj.get("type") != "user":
        return None
    msg = obj.get("message") or {}
    content = msg.get("content")
    if not isinstance(content, list):
        return None
    for item in content:
        if not isinstance(item, dict):
            continue
        src = item.get("source") or {}
        if item.get("type") == "image" and src.get("type") == "base64" and src.get("data"):
            try:
                data = base64.b64decode(src["data"])
            except Exception:
                continue
            return data, src.get("media_type", "image/png")
    return None


def _latest_pasted_image_from_log(log_path: str) -> tuple[bytes, str] | None:
    """扫描会话日志，返回最后一张粘贴图片 (bytes, mime)。"""
    result = None
    try:
        with open(log_path, "r", encoding="utf-8", errors="replace") as f:
            for line in f:
                found = _extract_pasted_image(line)
                if found:
                    result = found
    except OSError:
        return None
    return result


def _find_current_session_log() -> str | None:
    """定位当前会话的 JSONL 日志路径。

    优先用 CLAUDE_CODE_SESSION_ID 在 ~/.claude/projects 下按 mtime 找
    对应 `<sessionId>.jsonl`；找不到再退回"最近 10 分钟内修改的会话日志"。
    """
    projects = _claude_projects_dir()
    if not os.path.isdir(projects):
        return None
    sid = os.environ.get("CLAUDE_CODE_SESSION_ID") or os.environ.get("CLAUDE_SESSION_ID")
    now = time.time()

    def _recent_logs():
        for root, _dirs, files in os.walk(projects):
            for fn in files:
                if not fn.endswith(".jsonl"):
                    continue
                p = os.path.join(root, fn)
                try:
                    if now - os.path.getmtime(p) <= 600:
                        yield p
                except OSError:
                    continue

    if sid:
        for p in _recent_logs():
            if os.path.basename(p) == f"{sid}.jsonl":
                return p
        for root, _dirs, files in os.walk(projects):
            for fn in files:
                if fn == f"{sid}.jsonl":
                    return os.path.join(root, fn)
    try:
        return max(_recent_logs(), key=os.path.getmtime)
    except ValueError:
        return None


def _scan_recent_temp_images(seconds: int = 1800) -> list[str]:
    """兜底：在系统临时目录扫描近期图片文件（按文件签名识别），按 mtime 倒序。"""
    img_sign = re.compile(
        rb"\x89PNG\r\n\x1a\n|^\xff\xd8\xff|^RIFF.{4}WEBP|^\x00\x00\x00.f.t.y.p|^GIF8"
    )
    out = []
    for base in (tempfile.gettempdir(),):
        try:
            entries = os.scandir(base)
        except OSError:
            continue
        with entries as it:
            for e in it:
                try:
                    if not e.is_file():
                        continue
                    if time.time() - e.stat().st_mtime > seconds:
                        continue
                    with open(e.path, "rb") as f:
                        head = f.read(16)
                    if img_sign.match(head):
                        out.append(e.path)
                except OSError:
                    continue
    out.sort(key=os.path.getmtime, reverse=True)
    return out


def _resolve_pasted_image() -> tuple[bytes, str, str]:
    """找到用户最近粘贴的图片，返回 (bytes, mime, 来源描述)。

    顺序：当前会话日志 → 最近临时目录图片。找不到抛出 ValueError。
    """
    log = _find_current_session_log()
    if log:
        found = _latest_pasted_image_from_log(log)
        if found:
            data, mime = found
            return data, mime, f"会话日志 {os.path.basename(log)}"
    cands = _scan_recent_temp_images()
    if cands:
        path = cands[0]
        with open(path, "rb") as f:
            data = f.read()
        return data, _guess_mime(path, data), f"临时目录 {os.path.basename(path)}"
    raise ValueError("未找到粘贴的图片：既不在当前会话日志，临时目录也没有近期图片")


def _guess_mime(path: str, data: bytes) -> str:
    """按文件签名猜测图片 MIME（优先扩展名，其次魔数）。"""
    mime_by_ext = {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".webp": "image/webp", ".gif": "image/gif", ".avif": "image/avif", ".bmp": "image/bmp",
    }
    mime_by_sig = {
        b"\x89PNG": "image/png", b"\xff\xd8": "image/jpeg", b"RIFF": "image/webp",
        b"GIF8": "image/gif", b"ftyp": "image/avif",
    }
    ext = os.path.splitext(path)[1].lower()
    if ext in mime_by_ext:
        return mime_by_ext[ext]
    for sig, m in mime_by_sig.items():
        if data[:8].startswith(sig):
            return m
    return "image/jpeg"



def _find_soffice() -> str:
    for p in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]:
        if os.path.isfile(p):
            return p
    found = shutil.which("soffice")
    if found:
        return found
    raise RuntimeError("未找到 soffice.exe，请安装 LibreOffice")


def _find_pdftoppm() -> str:
    found = shutil.which("pdftoppm")
    if found:
        return found
    raise RuntimeError("未找到 pdftoppm，请安装 poppler")


def _ppt_page_to_png(pptx: str, page: int, dpi: int = 100) -> str:
    """把 PPTX 的指定页转成 PNG（soffice→pdf→pdftoppm），返回 PNG 临时文件路径。"""
    soffice = _find_soffice()
    pdftoppm = _find_pdftoppm()
    work = tempfile.mkdtemp(prefix="ppt_")
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", work, pptx],
            check=True, timeout=180,
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        pdf = os.path.join(work, os.path.splitext(os.path.basename(pptx))[0] + ".pdf")
        if not os.path.isfile(pdf):
            raise RuntimeError("soffice 未生成 PDF")
        out_png = os.path.join(work, "page")
        subprocess.run(
            [pdftoppm, "-png", "-r", str(dpi), "-f", str(page), "-l", str(page), pdf, out_png],
            check=True, timeout=60,
        )
        # poppler 输出命名可能带前导零（page-01.png），用通配符匹配
        import glob
        candidates = glob.glob(f"{out_png}-*.png")
        if not candidates:
            raise RuntimeError("pdftoppm 未生成 PNG")
        return candidates[0]
    except Exception:
        shutil.rmtree(work, ignore_errors=True)
        raise


@mcp.tool()
def analyze_ppt(
    pptx: str,
    page: int | None = None,
    prompt: str = "请用中文描述这一页 PPT 的内容和排版特点（布局、配色、图文结构）。",
    dpi: int = 100,
) -> str:
    """把 PPT 的某一页（或全部分页）渲染成图后交给视觉模型分析（适合看排版效果）。

    Args:
        pptx: .pptx 文件路径。
        page: 要分析的页码（从 1 开始）。不传则分析全部页，逐页返回。
        prompt: 想要模型关注什么（默认描述内容与排版）。
        dpi: 渲染分辨率，默认 100。
    """
    if not os.path.isfile(pptx):
        return f"[vision-mcp 错误] 文件不存在: {pptx}"
    try:
        if page is not None:
            png = _ppt_page_to_png(pptx, page, dpi)
            try:
                return _analyze(png, prompt)
            finally:
                if os.path.exists(png):
                    os.remove(png)
        else:
            # 全部页：先 soffice 转 pdf，再逐页 pdftoppm + _analyze
            soffice = _find_soffice()
            pdftoppm = _find_pdftoppm()
            work = tempfile.mkdtemp(prefix="ppt_")
            try:
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf", "--outdir", work, pptx],
                    check=True, timeout=180,
                    stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                )
                pdf = os.path.join(work, os.path.splitext(os.path.basename(pptx))[0] + ".pdf")
                if not os.path.isfile(pdf):
                    raise RuntimeError("soffice 未生成 PDF")
                out_png = os.path.join(work, "page")
                subprocess.run(
                    [pdftoppm, "-png", "-r", str(dpi), pdf, out_png],
                    check=True, timeout=120,
                )
                import glob
                pages = sorted(glob.glob(f"{out_png}-*.png"))
                results = []
                for i, p in enumerate(pages, 1):
                    try:
                        txt = _analyze(p, f"{prompt}（第{i}页）")
                        results.append(f"=== 第 {i} 页 ===\n{txt}")
                    finally:
                        if os.path.exists(p):
                            os.remove(p)
                    # 页间间隔，避免触发火山限流（多页连续请求易 429）
                    if i < len(pages):
                        time.sleep(4)
                return "\n\n".join(results) if results else "[vision-mcp 错误] 未生成任何页面"
            finally:
                shutil.rmtree(work, ignore_errors=True)
    except Exception as e:
        return f"[vision-mcp 错误] {e}"


def _shape_rect(sh) -> tuple[float, float, float, float] | None:
    """返回形状包围盒 (l, t, r, b)，单位 pt；无尺寸返回 None。"""
    if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
        return None
    EMU = 12700.0
    l, t = sh.left / EMU, sh.top / EMU
    return l, t, l + sh.width / EMU, t + sh.height / EMU


def _rect_overlap_area(a: tuple, b: tuple) -> float:
    l = max(a[0], b[0]); t = max(a[1], b[1])
    r = min(a[2], b[2]); btm = min(a[3], b[3])
    w, h = r - l, btm - t
    return w * h if w > 0 and h > 0 else 0.0


def _shape_label(sh) -> str:
    txt = ""
    if sh.has_text_frame and sh.text_frame.text.strip():
        txt = sh.text_frame.text.strip().replace("\n", " ")[:20]
    return f"{sh.shape_type}[{txt}]"


def _ppt_overlaps(pptx: str, min_overlap_pt: float = 2.0) -> list[dict]:
    """用 python-pptx 扫描所有页，返回包围盒重叠的形状对列表。

    每项: {page, a, b, area_pt2}
    """
    from pptx import Presentation
    prs = Presentation(pptx)
    results = []
    thresh = min_overlap_pt * min_overlap_pt
    for si, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        for i in range(len(shapes)):
            for j in range(i + 1, len(shapes)):
                ra = _shape_rect(shapes[i])
                rb = _shape_rect(shapes[j])
                if ra is None or rb is None:
                    continue
                area = _rect_overlap_area(ra, rb)
                if area > thresh:
                    results.append({
                        "page": si,
                        "a": _shape_label(shapes[i]),
                        "b": _shape_label(shapes[j]),
                        "area_pt2": round(area, 1),
                    })
    return results


@mcp.tool()
def check_ppt_overlap(
    pptx: str,
    min_overlap_pt: float = 2.0,
    visual_check: bool = True,
) -> str:
    """检测 PPT 排版元素是否重叠（坐标检测 + 视觉交叉验证双保险）。

    先精确读每个形状的包围盒坐标，找出矩形重叠超过阈值的对；
    再用视觉模型看渲染图，确认这些重叠在视觉上是否真的可见
    （文本框包围盒重叠但文字有内边距时，视觉上可能并不重叠）。

    Args:
        pptx: .pptx 文件路径。
        min_overlap_pt: 视为重叠的最小重叠边长（pt），默认 2。
        visual_check: 是否做视觉交叉验证，默认 True。
    """
    if not os.path.isfile(pptx):
        return f"[vision-mcp 错误] 文件不存在: {pptx}"
    try:
        overlaps = _ppt_overlaps(pptx, min_overlap_pt)
        if not overlaps:
            return "未检测到任何坐标级重叠。\n\n（python-pptx 扫描所有页，包围盒无相交。）"

        # 按页分组
        by_page = {}
        for o in overlaps:
            by_page.setdefault(o["page"], []).append(o)

        lines = [f"坐标检测到 {len(overlaps)} 处包围盒重叠（阈值 {min_overlap_pt}pt）：", ""]
        for pg in sorted(by_page):
            lines.append(f"第 {pg} 页:")
            for o in by_page[pg]:
                lines.append(f"  · {o['a']} 与 {o['b']} 重叠 {o['area_pt2']} pt²")
            lines.append("")

        if not visual_check:
            return "\n".join(lines)

        # 视觉交叉验证：逐页渲染确认
        lines.append("=== 视觉交叉验证 ===")
        lines.append("")
        vis_pages = sorted(by_page)
        for idx, pg in enumerate(vis_pages):
            png = _ppt_page_to_png(pptx, pg)
            try:
                pairs = "；".join(
                    f"{o['a']} 与 {o['b']}" for o in by_page[pg]
                )
                vis = _analyze(
                    png,
                    "这一页 PPT 上，是否有元素在视觉上确实重叠/遮盖？"
                    f"（坐标为：{pairs}）"
                    " 请只回答：确实重叠 / 不重叠（是包围盒重叠但文字有间距）。",
                )
                lines.append(f"第 {pg} 页 → {vis}")
                lines.append("")
            finally:
                if os.path.exists(png):
                    os.remove(png)
            if idx < len(vis_pages) - 1:
                time.sleep(4)

        return "\n".join(lines)
    except Exception as e:
        return f"[vision-mcp 错误] {e}"


@mcp.tool()
def analyze_image(
    image: str,
    prompt: str = "请详细描述这张图片的内容。",
    detail: bool = False,
    thinking: bool | None = None,
) -> str:
    """分析一张图片，返回视觉模型的文字描述。

    Args:
        image: 图片来源 —— 本地文件路径，或 http(s) URL，或 data URL，或 base64。
        prompt: 想要模型关注什么（默认详细描述）。
        detail: 设为 True 时启用完整描述模式 —— 结构化分区 + 查漏补缺，
                尽量捕获全部文字/数字/表格，适合学术文献页等需要完整信息的场景。
        thinking: None=按场景默认（detail 开、其余关）；True/False=强制覆盖。
    """
    try:
        data, mime, _ = _read_image_bytes(image)
        if len(data) > MAX_IMAGE_BYTES:
            return f"[vision-mcp 错误] 图片过大（{len(data)} bytes），上限 {MAX_IMAGE_BYTES}"
        fd, tmp = tempfile.mkstemp(suffix=".png")
        try:
            with os.fdopen(fd, "wb") as f:
                f.write(data)
            think = _resolve_thinking(detail, thinking)
            if detail:
                return _analyze_detailed(tmp, prompt, thinking=think)
            return _analyze(tmp, prompt, thinking=think)
        finally:
            if os.path.exists(tmp):
                os.remove(tmp)
    except Exception as e:
        return f"[vision-mcp 错误] {e}"


@mcp.tool()
def describe_image(image: str) -> str:
    """快速浏览图片：输出一段通用描述（适合不指定任务的场合）。"""
    try:
        return _analyze(image, "请用简洁的中文描述这张图片的大致内容。", thinking=False)
    except Exception as e:
        return f"[vision-mcp 错误] {e}"


@mcp.tool()
def analyze_pasted_image(
    prompt: str = "请用中文详细描述这张图片的内容。",
    detail: bool = False,
    thinking: bool | None = None,
) -> str:
    """分析对话中粘贴的最近一张图片（无需提供文件路径）。

    从当前 Claude Code 会话日志里提取用户粘贴的最后一张图片交给视觉模型，
    找不到时回退到系统临时目录扫描近期图片文件。

    Args:
        prompt: 想要模型关注什么（默认详细描述）。
        detail: 设为 True 时启用完整描述模式 —— 结构化分区 + 查漏补缺，
                尽量捕获全部文字/数字/表格，适合学术文献页等需要完整信息的场景。
        thinking: None=按场景默认（detail 开、其余关）；True/False=强制覆盖。
    """
    try:
        data, mime, source = _resolve_pasted_image()
        if len(data) > MAX_IMAGE_BYTES:
            return f"[vision-mcp 错误] 图片过大（{len(data)} bytes），上限 {MAX_IMAGE_BYTES}"
        fd, tmp = tempfile.mkstemp(suffix=".png")
        try:
            with os.fdopen(fd, "wb") as f:
                f.write(data)
            think = _resolve_thinking(detail, thinking)
            if detail:
                result = _analyze_detailed(tmp, prompt, thinking=think)
            else:
                result = _analyze(tmp, prompt, thinking=think)
            return result + f"\n\n（图片来源：{source}）"
        finally:
            if os.path.exists(tmp):
                os.remove(tmp)
    except Exception as e:
        return f"[vision-mcp 错误] {e}"


@mcp.tool()
def analyze_html(
    url: str,
    prompt: str = "请用中文描述这个网页的布局、视觉风格和主要内容。",
    width: int = 1280,
    height: int = 900,
) -> str:
    """把 HTML 渲染成截图后交给视觉模型分析（适合看网页视觉效果）。

    Args:
        url: 网页 URL，或本地 HTML 文件路径，或 data:text/html,... 内容。
        prompt: 想要模型关注什么（默认描述布局/风格/内容）。
        width: 截图宽度，默认 1280。
        height: 截图高度，默认 900。
    """
    try:
        png = _html_to_png(url, width, height)
        try:
            return _analyze(png, prompt)
        finally:
            if os.path.exists(png):
                os.remove(png)
    except Exception as e:
        return f"[vision-mcp 错误] {e}"


if __name__ == "__main__":
    mcp.run(transport="stdio")
